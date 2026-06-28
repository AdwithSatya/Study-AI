"""
ai/ingest.py
────────────────────────────────────────────────────────────────────────────────
Document ingestion pipeline.

Responsibilities:
  1. Parse raw file bytes into plain text (PDF, PPTX, DOCX, TXT, MD).
  2. Split text into overlapping chunks for better retrieval coverage.
  3. Embed each chunk via the Hugging Face Inference API (no local model).
  4. Store chunks + embeddings + metadata in ChromaDB.

Entry point:
    chunks_stored = ingest(file_bytes, filename, file_id, folder_id, user_id)

Each chunk is stored with metadata so it can be filtered per-user and per-folder
at query time, keeping workspaces isolated from each other.

NOTE: Embeddings are generated through the HuggingFace Inference API
      (HuggingFaceEndpointEmbeddings) — no model weights are downloaded,
      keeping the Render free-tier deployment well under the 525 MB limit.
"""

import os
from dotenv import load_dotenv
from langchain_huggingface.embeddings import HuggingFaceEndpointEmbeddings
from langchain_text_splitters import RecursiveCharacterTextSplitter
import chromadb

load_dotenv()


def _require_env(name: str) -> str:
    """Read a required environment variable; raise clearly if it is missing."""
    value = os.getenv(name)
    if value is None:
        raise RuntimeError(
            f"Required environment variable '{name}' is not set. "
            "Check your .env file."
        )
    return value


# ── Embedding model (remote API — zero local disk footprint) ───────────────────
_model_name = os.getenv(
    "EMBEDDING_MODEL_NAME", "sentence-transformers/all-MiniLM-L6-v2"
)
embedding_model = HuggingFaceEndpointEmbeddings(
    model=_model_name,
    huggingfacehub_api_token=_require_env("HF_API_TOKEN"),
)

# ── ChromaDB cloud client ──────────────────────────────────────────────────────
_chroma_client = chromadb.CloudClient(
    api_key=_require_env("CHROMA_API_KEY"),
    tenant=_require_env("CHROMA_TENANT"),
    database=_require_env("CHROMA_DATABASE"),
)
collection = _chroma_client.get_or_create_collection(name="study_notes")

# ── Text splitter (shared config) ─────────────────────────────────────────────
_splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=100)

SUPPORTED_EXTENSIONS = {"pdf", "ppt", "pptx", "docx", "txt", "md"}


def _extract_text(file_bytes: bytes, filename: str) -> str:
    """Extract plain text from a supported file type."""
    ext = filename.rsplit(".", 1)[-1].lower()

    if ext == "pdf":
        import fitz  # PyMuPDF
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        return "\n".join(
            page.get_text().strip() for page in doc if page.get_text().strip()
        )

    if ext in {"ppt", "pptx"}:
        from pptx import Presentation
        from io import BytesIO
        prs = Presentation(BytesIO(file_bytes))
        return "\n".join(
            shape.text
            for slide in prs.slides
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        )

    if ext == "docx":
        import docx
        from io import BytesIO
        document = docx.Document(BytesIO(file_bytes))
        return "\n".join(p.text for p in document.paragraphs if p.text.strip())

    if ext in {"txt", "md"}:
        return file_bytes.decode("utf-8", errors="replace")

    raise ValueError(f"Unsupported file type: .{ext}")


def ingest(
    file_bytes: bytes,
    filename: str,
    file_id: str,
    folder_id: str,
    user_id: str,
) -> int:
    """
    Parse, chunk, embed, and store a document in ChromaDB.

    Embeddings are generated via the HuggingFace Inference API — no local
    model weights are loaded, so memory usage stays within Render's free tier.

    Returns:
        Number of chunks stored.
    """
    text = _extract_text(file_bytes, filename)
    if not text.strip():
        raise ValueError("Document appears to be empty or unreadable.")

    chunks = _splitter.split_text(text)

    # embed_documents batches all chunks in one API call
    embeddings = embedding_model.embed_documents(chunks)

    for i, (chunk, embedding) in enumerate(zip(chunks, embeddings)):
        collection.add(
            ids=[f"{file_id}_{i}"],
            embeddings=[embedding],
            documents=[chunk],
            metadatas=[{
                "user_id": user_id,
                "folder_id": folder_id,
                "file_id": file_id,
                "filename": filename,
            }],
        )

    return len(chunks)


def delete_embeddings(file_id: str | None = None, folder_id: str | None = None) -> None:
    """
    Remove vectors from ChromaDB.

    Pass file_id to delete a single file's chunks.
    Pass folder_id to wipe an entire workspace.
    """
    if file_id:
        collection.delete(where={"file_id": file_id})
    elif folder_id:
        collection.delete(where={"folder_id": folder_id})
